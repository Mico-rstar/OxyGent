"""
文档读取工具 - 支持txt, pdf, ppt, xlsx文档内容读取
基于阿里云通义千问OCR模型进行PDF和PPT的内容解析
"""

import os
import json
import tempfile
from typing import Optional, Dict, Any
from pydantic import Field
from pathlib import Path

import fitz  # PyMuPDF
import pandas as pd
from dashscope import MultiModalConversation
from pptx import Presentation
import pyarrow.parquet as pq

from oxygent.oxy import FunctionHub

document_tools = FunctionHub(name="document_tools")


def _convert_pdf_to_images(pdf_path: str, output_dir: str) -> list:
    """
    将PDF转换为图片，用于OCR处理

    Args:
        pdf_path: PDF文件路径
        output_dir: 输出目录

    Returns:
        图片文件路径列表
    """

    doc = fitz.open(pdf_path)
    image_paths = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        mat = fitz.Matrix(2.0, 2.0)  # 2x缩放提高清晰度
        pix = page.get_pixmap(matrix=mat)

        image_path = os.path.join(output_dir, f"page_{page_num + 1}.png")
        pix.save(image_path)
        image_paths.append(image_path)

    doc.close()
    return image_paths




@document_tools.tool(description="读取TXT文档内容")
def read_txt_document(
    filename: str = Field(description="TXT文档的文件名"),
    encoding: str = Field(description="文件编码，如utf-8")
) -> str:
    """
    读取TXT文档内容

    Args:
        file_path: TXT文档文件的绝对路径
        encoding: 文件编码

    Returns:
        文档内容字符串

    Raises:
        FileNotFoundError: 如果文件不存在
        UnicodeDecodeError: 如果编码错误
    """
    file_path = os.getenv("DEFAULT_FILE_PATH") + filename
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    try:
        with open(file_path, 'r', encoding=encoding) as f:
            content = f.read()
        return content
    except UnicodeDecodeError:
        # 尝试其他编码
        for enc in ['gbk', 'gb2312', 'latin-1']:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    content = f.read()
                return f"使用编码 {enc} 读取:\n\n{content}"
            except UnicodeDecodeError:
                continue
        raise UnicodeDecodeError(f"无法解码文件 {file_path}")


@document_tools.tool(description="读取PDF文档内容")
def read_pdf_document(
    filename: str = Field(description="PDF文档文件的文件名"),
    question: str = Field(description="对文档的特定要求"),
) -> str:
    """
    使用OCR读取PDF文档内容

    Args:
        file_path: PDF文档文件的绝对路径
        question: 对文档的特定要求

    Returns:
        文档内容的文本描述

    Raises:
        ImportError: 如果缺少依赖库
        FileNotFoundError: 如果文件不存在
        ValueError: 如果API调用失败
    """

    file_path = os.getenv("DEFAULT_FILE_PATH") + filename
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PDF文件不存在: {file_path}")

    # 使用提供的API Key或环境变量
    key = os.getenv('DASHSCOPE_API_KEY')
    if not key:
        raise ValueError("未找到DashScope API Key，请设置环境变量DASHSCOPE_API_KEY")

    try:
        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            # 将PDF转换为图片
            image_paths = _convert_pdf_to_images(file_path, temp_dir)
            all_content = []
            model = "qwen-vl-ocr-latest"

            # 处理每一页
            for i, image_path in enumerate(image_paths):
                try:
                    print(image_path, i)
                    # 构建消息
                    messages = [
                        {
                            'role': 'user',
                            'content': [
                                {'image': f"file://{image_path}"},
                                {'text': question}
                            ]
                        }
                    ]
                    # 调用OCR API
                    response = MultiModalConversation.call(
                        api_key=key,
                        model=model,
                        messages=messages
                    )
                    print("complish")
                    if response.status_code == 200:
                        page_content = response.output.choices[0].message.content[0]['text']
                        all_content.append(f"=== 第 {i + 1} 页 ===\n{page_content}")
                    else:
                        all_content.append(f"=== 第 {i + 1} 页 ===\nOCR处理失败: {response.message}")

                except Exception as e:
                    all_content.append(f"=== 第 {i + 1} 页 ===\n处理错误: {str(e)}")

            return "\n\n".join(all_content)

    except Exception as e:
        raise ValueError(f"PDF文档读取失败: {str(e)}")


@document_tools.tool(description="读取PPT文档内容")
def read_ppt_document(
    filename: str = Field(description="PPT文档的文件名"),
) -> str:
    """
    读取PPT文档内容

    Args:
        file_path: PPT文档文件的绝对路径

    Returns:
        文档内容的文本描述

    Raises:
        ImportError: 如果缺少依赖库
        FileNotFoundError: 如果文件不存在
    """
    file_path = os.getenv("DEFAULT_FILE_PATH") + filename
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPT文件不存在: {file_path}")

    try:
        prs = Presentation(file_path)
        all_content = []

        # 读取每张幻灯片的内容
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_content.append(f"=== 第 {slide_num} 张幻灯片 ===")

            # 读取幻灯片标题和内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

            # 如果幻灯片有内容，添加到结果中
            if len(slide_content) > 1:  # 除了标题行还有内容
                all_content.append("\n".join(slide_content))
            else:
                all_content.append(f"=== 第 {slide_num} 张幻灯片 ===\n（空白幻灯片或无文本内容）")

        if not all_content:
            return "PPT文档中没有找到文本内容"

        return "\n\n".join(all_content)

    except Exception as e:
        raise ValueError(f"PPT文档读取失败: {str(e)}")



@document_tools.tool(description="读取Parquet文档内容")
def read_parquet_document(
    filename: str = Field(description="Parquet文档的文件名"),
    max_rows: int = Field(default=1000, description="最大读取行数，默认1000行"),
) -> str:
    """
    读取Parquet文档内容

    Args:
        file_path: Parquet文档文件的绝对路径
        max_rows: 最大读取行数，默认1000行

    Returns:
        文档内容的格式化字符串

    Raises:
        ImportError: 如果缺少依赖库
        FileNotFoundError: 如果文件不存在
        ValueError: 如果API调用失败
    """
    file_path = os.getenv("DEFAULT_FILE_PATH") + filename
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Parquet文件不存在: {file_path}")

    try:
        # 读取Parquet文件的元数据
        parquet_file = pq.ParquetFile(file_path)

        # 获取schema信息
        schema = parquet_file.schema_arrow
        num_rows = parquet_file.metadata.num_rows
        num_columns = len(schema)

        result_parts = []
        result_parts.append(f"=== Parquet文件信息 ===")
        result_parts.append(f"文件路径: {file_path}")
        result_parts.append(f"总行数: {num_rows:,}")
        result_parts.append(f"总列数: {num_columns}")
        result_parts.append(f"列信息: {[field.name for field in schema]}")
        result_parts.append("")

        # 使用pandas读取数据（限制行数以避免内存问题）
        df = pd.read_parquet(file_path)

        if df.empty:
            result_parts.append("Parquet文件为空")
            return "\n".join(result_parts)

        # 如果数据量很大，只显示前几行
        if len(df) > max_rows:
            result_parts.append(f"数据预览（前 {max_rows} 行，共 {len(df):,} 行）:")
            df_preview = df.head(max_rows)
        else:
            result_parts.append(f"完整数据（共 {len(df):,} 行）:")
            df_preview = df

        # 转换为字符串格式
        result_parts.append(f"{df_preview.to_string(index=False)}")

        # 添加数据类型信息
        result_parts.append("")
        result_parts.append("=== 数据类型信息 ===")
        for col, dtype in df.dtypes.items():
            result_parts.append(f"{col}: {dtype}")

        # 添加基本统计信息（仅对数值列）
        numeric_cols = df.select_dtypes(include=['number']).columns
        if not numeric_cols.empty:
            result_parts.append("")
            result_parts.append("=== 数值列统计信息 ===")
            stats = df[numeric_cols].describe()
            result_parts.append(f"{stats.to_string()}")

        return "\n".join(result_parts)

    except Exception as e:
        raise ValueError(f"Parquet文档读取失败: {str(e)}")


@document_tools.tool(description="读取Excel文档内容")
def read_xlsx_document(
    filename: str = Field(description="Excel文档的文件名"),
) -> str:
    """
    读取Excel文档内容

    Args:
        file_path: Excel文档文件的绝对路径
        sheet_name: 工作表名称，None表示读取所有工作表
        max_rows: 每个工作表最大读取行数

    Returns:
        文档内容的格式化字符串

    Raises:
        ImportError: 如果缺少依赖库
        FileNotFoundError: 如果文件不存在
    """
    file_path = os.getenv("DEFAULT_FILE_PATH") + filename
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Excel文件不存在: {file_path}")

    try:
        # 读取所有工作表
        excel_file = pd.ExcelFile(file_path)
        result_parts = []

        for sheet in excel_file.sheet_names:
            try:
                df = pd.read_excel(file_path, sheet_name=sheet, nrows=1000)
                if not df.empty:
                    result_parts.append(f"=== 工作表: {sheet} ===\n{df.to_string(index=False)}")
            except Exception as e:
                result_parts.append(f"=== 工作表: {sheet} ===\n读取失败: {str(e)}")

        result = "\n\n".join(result_parts)

        return result

    except Exception as e:
        raise ValueError(f"Excel文档读取失败: {str(e)}")




